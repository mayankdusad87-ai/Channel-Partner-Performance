from pptx import Presentation
from datetime import datetime

def create_ppt(insights, summary):

    prs = Presentation()

    # TITLE
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Channel Partner Strategy Deck"
    slide.placeholders[1].text = str(datetime.now().date())

    # EXEC SUMMARY
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    slide.placeholders[1].text = insights

    # NETWORK SPLIT
    strat = summary["Strategy"].value_counts()

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "CP Strategy Split"
    slide.placeholders[1].text = strat.to_string()

    # TOP CP
    top = summary.sort_values("Bookings", ascending=False).head(5)

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Top Performers"
    slide.placeholders[1].text = top.to_string(index=False)

    # ACTION TABLE
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Action Plan"
    slide.placeholders[1].text = summary[["Channel Partner","Strategy","Action"]].to_string(index=False)

    file = "Strategy_Report.pptx"
    prs.save(file)

    return file
